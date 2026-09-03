"""Persistent data and security helpers for the ATLAS faculty pilot.

The pilot is intentionally small: at most five invited Northeastern faculty
members, one Railway application replica, and one SQLite database stored on a
persistent Railway volume. Course documents live beside the database on that
same volume. Each professor can supply separate Anthropic and OpenAI API keys,
which are encrypted before they are written to disk.
"""

from __future__ import annotations

import base64
import hashlib
import hmac
import json
import os
import re
import secrets
import sqlite3
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from cryptography.fernet import Fernet, InvalidToken

from ai_providers import DEFAULT_MODEL_BY_PROVIDER, normalize_provider, validate_provider_model


DEFAULT_MAX_PROFESSORS = 5
DEFAULT_MONTHLY_QUESTION_LIMIT = 500
MAX_DOCUMENT_BYTES = 25 * 1024 * 1024
MAX_COURSE_DOCUMENT_BYTES = 250 * 1024 * 1024
MAX_EXTRACTED_TEXT_CHARS = 1_000_000
MAX_COURSE_EXTRACTED_CHARS = 5_000_000
MAX_COURSE_DOCUMENTS = 100
MAX_COURSES_PER_PROFESSOR = 10
ALLOWED_DOCUMENT_EXTENSIONS = {".docx", ".md", ".pdf", ".pptx", ".txt"}
NORTHEASTERN_EMAIL_PATTERN = re.compile(
    r"^[A-Z0-9._%+\-]+@northeastern\.edu$", re.IGNORECASE
)


class PilotConfigurationError(RuntimeError):
    """Raised when the enabled pilot is missing a required secret or path."""


class PilotValidationError(ValueError):
    """Raised for a safe, user-facing pilot validation failure."""


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def normalize_email(email: str) -> str:
    normalized = (email or "").strip().lower()
    if not NORTHEASTERN_EMAIL_PATTERN.fullmatch(normalized):
        raise PilotValidationError("A valid @northeastern.edu email is required.")
    return normalized


def hash_token(token: str) -> str:
    return hashlib.sha256(token.encode("utf-8")).hexdigest()


def hash_password(password: str) -> str:
    """Hash a faculty password with the standard-library scrypt KDF."""
    if len(password or "") < 12:
        raise PilotValidationError("The password must contain at least 12 characters.")
    salt = secrets.token_bytes(16)
    digest = hashlib.scrypt(
        password.encode("utf-8"), salt=salt, n=2**14, r=8, p=1, dklen=32
    )
    return "scrypt$16384$8$1${}${}".format(
        base64.urlsafe_b64encode(salt).decode("ascii"),
        base64.urlsafe_b64encode(digest).decode("ascii"),
    )


def verify_password(password: str, stored_hash: str) -> bool:
    try:
        scheme, n, r, p, salt_b64, digest_b64 = stored_hash.split("$", 5)
        if scheme != "scrypt":
            return False
        salt = base64.urlsafe_b64decode(salt_b64.encode("ascii"))
        expected = base64.urlsafe_b64decode(digest_b64.encode("ascii"))
        actual = hashlib.scrypt(
            password.encode("utf-8"),
            salt=salt,
            n=int(n),
            r=int(r),
            p=int(p),
            dklen=len(expected),
        )
        return hmac.compare_digest(actual, expected)
    except (ValueError, TypeError):
        return False


_DUMMY_PASSWORD_HASH = hash_password("invalid-account-password")


def generate_encryption_key() -> str:
    """Return a new Fernet key suitable for ATLAS_ENCRYPTION_KEY."""
    return Fernet.generate_key().decode("ascii")


def _clean_text(value: str, label: str, max_length: int, required: bool = True) -> str:
    cleaned = " ".join((value or "").split())
    if required and not cleaned:
        raise PilotValidationError(f"{label} is required.")
    if len(cleaned) > max_length:
        raise PilotValidationError(f"{label} must be {max_length} characters or fewer.")
    return cleaned


def safe_filename(filename: str) -> str:
    name = Path(filename or "").name.strip()
    if not name or name in {".", ".."}:
        raise PilotValidationError("A valid filename is required.")
    if len(name) > 180:
        raise PilotValidationError("The filename is too long.")
    extension = Path(name).suffix.lower()
    if extension not in ALLOWED_DOCUMENT_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_DOCUMENT_EXTENSIONS))
        raise PilotValidationError(f"Supported file types are: {allowed}.")
    return name


def extract_document_text(filename: str, content: bytes) -> str:
    """Extract searchable text from a supported faculty document."""
    if len(content) > MAX_DOCUMENT_BYTES:
        raise PilotValidationError("Each uploaded document must be 25 MB or smaller.")

    name = safe_filename(filename)
    extension = Path(name).suffix.lower()

    text = ""
    try:
        if extension in {".md", ".txt"}:
            text = content.decode("utf-8", errors="replace")

        elif extension == ".docx":
            from io import BytesIO
            from docx import Document

            document = Document(BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)

        elif extension == ".pdf":
            from io import BytesIO
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(content))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)

        elif extension == ".pptx":
            from io import BytesIO
            from pptx import Presentation

            presentation = Presentation(BytesIO(content))
            lines: List[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text)
            text = "\n".join(lines)
    except Exception as exc:
        raise PilotValidationError(f"ATLAS could not read {name}: {exc}") from exc

    if not text and extension not in ALLOWED_DOCUMENT_EXTENSIONS:
        raise PilotValidationError("Unsupported document type.")
    if len(text) > MAX_EXTRACTED_TEXT_CHARS:
        raise PilotValidationError(
            "The extracted document text is too large. Split the file into smaller parts."
        )
    return text


class PilotStore:
    """SQLite and filesystem-backed persistence for the faculty pilot."""

    def __init__(
        self,
        data_dir: Path,
        encryption_key: str,
        max_professors: int = DEFAULT_MAX_PROFESSORS,
    ) -> None:
        self.data_dir = Path(data_dir).resolve()
        self.data_dir.mkdir(parents=True, exist_ok=True)
        try:
            self.data_dir.chmod(0o700)
        except OSError:
            pass
        self.courses_dir = self.data_dir / "courses"
        self.courses_dir.mkdir(parents=True, exist_ok=True)
        try:
            self.courses_dir.chmod(0o700)
        except OSError:
            pass
        self.db_path = self.data_dir / "atlas_pilot.sqlite3"
        self.max_professors = max(1, int(max_professors))
        try:
            self.fernet = Fernet(encryption_key.encode("ascii"))
        except (ValueError, TypeError) as exc:
            raise PilotConfigurationError(
                "ATLAS_ENCRYPTION_KEY must be a valid Fernet key."
            ) from exc

    def _connect(self) -> sqlite3.Connection:
        connection = sqlite3.connect(self.db_path, timeout=30)
        connection.row_factory = sqlite3.Row
        connection.execute("PRAGMA foreign_keys = ON")
        connection.execute("PRAGMA journal_mode = WAL")
        return connection

    def initialize(self) -> None:
        with self._connect() as connection:
            connection.executescript(
                """
                CREATE TABLE IF NOT EXISTS professors (
                    id TEXT PRIMARY KEY,
                    email TEXT NOT NULL UNIQUE COLLATE NOCASE,
                    name TEXT NOT NULL,
                    password_hash TEXT NOT NULL,
                    api_key_encrypted TEXT,
                    api_key_last_four TEXT,
                    openai_api_key_encrypted TEXT,
                    openai_api_key_last_four TEXT,
                    is_active INTEGER NOT NULL DEFAULT 1,
                    created_at TEXT NOT NULL
                );

                CREATE TABLE IF NOT EXISTS invitations (
                    id TEXT PRIMARY KEY,
                    email TEXT NOT NULL COLLATE NOCASE,
                    name TEXT NOT NULL,
                    token_hash TEXT NOT NULL UNIQUE,
                    expires_at TEXT NOT NULL,
                    used_at TEXT,
                    created_at TEXT NOT NULL
                );

                CREATE TABLE IF NOT EXISTS courses (
                    id TEXT PRIMARY KEY,
                    owner_id TEXT NOT NULL,
                    name TEXT NOT NULL,
                    code TEXT NOT NULL,
                    section TEXT NOT NULL DEFAULT '',
                    term TEXT NOT NULL,
                    campus TEXT NOT NULL DEFAULT 'Arlington',
                    status TEXT NOT NULL DEFAULT 'draft',
                    monthly_question_limit INTEGER NOT NULL DEFAULT 500,
                    concept_map_json TEXT NOT NULL DEFAULT '{}',
                    provider TEXT NOT NULL DEFAULT 'anthropic',
                    model TEXT NOT NULL DEFAULT 'claude-sonnet-4-6',
                    model_updated_at TEXT,
                    project_builder_enabled INTEGER NOT NULL DEFAULT 1,
                    research_innovation_enabled INTEGER NOT NULL DEFAULT 0,
                    created_at TEXT NOT NULL,
                    updated_at TEXT NOT NULL,
                    FOREIGN KEY (owner_id) REFERENCES professors(id)
                );

                CREATE TABLE IF NOT EXISTS sessions (
                    token_hash TEXT PRIMARY KEY,
                    role TEXT NOT NULL,
                    professor_id TEXT,
                    expires_at TEXT NOT NULL,
                    created_at TEXT NOT NULL,
                    FOREIGN KEY (professor_id) REFERENCES professors(id) ON DELETE CASCADE
                );

                CREATE TABLE IF NOT EXISTS documents (
                    id TEXT PRIMARY KEY,
                    course_id TEXT NOT NULL,
                    filename TEXT NOT NULL,
                    document_type TEXT NOT NULL,
                    stored_path TEXT NOT NULL,
                    extracted_path TEXT NOT NULL,
                    byte_size INTEGER NOT NULL,
                    extracted_chars INTEGER NOT NULL DEFAULT 0,
                    sha256 TEXT NOT NULL,
                    uploaded_at TEXT NOT NULL,
                    FOREIGN KEY (course_id) REFERENCES courses(id) ON DELETE CASCADE
                );

                CREATE TABLE IF NOT EXISTS usage_events (
                    id TEXT PRIMARY KEY,
                    course_id TEXT NOT NULL,
                    professor_id TEXT NOT NULL,
                    session_id TEXT NOT NULL,
                    provider TEXT NOT NULL DEFAULT 'anthropic',
                    model TEXT NOT NULL,
                    event_type TEXT NOT NULL DEFAULT 'student_question',
                    input_tokens INTEGER NOT NULL,
                    output_tokens INTEGER NOT NULL,
                    created_at TEXT NOT NULL,
                    FOREIGN KEY (course_id) REFERENCES courses(id) ON DELETE CASCADE,
                    FOREIGN KEY (professor_id) REFERENCES professors(id)
                );

                CREATE TABLE IF NOT EXISTS pilot_feedback (
                    id TEXT PRIMARY KEY,
                    course_id TEXT NOT NULL,
                    session_id TEXT NOT NULL,
                    rating TEXT NOT NULL,
                    comment TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL,
                    FOREIGN KEY (course_id) REFERENCES courses(id) ON DELETE CASCADE
                );

                CREATE TABLE IF NOT EXISTS course_model_events (
                    id TEXT PRIMARY KEY,
                    course_id TEXT NOT NULL,
                    professor_id TEXT NOT NULL,
                    old_provider TEXT,
                    old_model TEXT,
                    new_provider TEXT NOT NULL,
                    new_model TEXT NOT NULL,
                    created_at TEXT NOT NULL,
                    FOREIGN KEY (course_id) REFERENCES courses(id) ON DELETE CASCADE,
                    FOREIGN KEY (professor_id) REFERENCES professors(id)
                );

                CREATE TABLE IF NOT EXISTS legacy_course_features (
                    course_id TEXT PRIMARY KEY,
                    project_builder_enabled INTEGER NOT NULL DEFAULT 1,
                    research_innovation_enabled INTEGER NOT NULL DEFAULT 0,
                    updated_at TEXT NOT NULL
                );

                CREATE INDEX IF NOT EXISTS idx_courses_owner
                    ON courses(owner_id);
                CREATE INDEX IF NOT EXISTS idx_sessions_expires
                    ON sessions(expires_at);
                CREATE INDEX IF NOT EXISTS idx_documents_course
                    ON documents(course_id);
                CREATE INDEX IF NOT EXISTS idx_usage_course_created
                    ON usage_events(course_id, created_at);
                CREATE INDEX IF NOT EXISTS idx_course_model_events
                    ON course_model_events(course_id, created_at);
                """
            )
            professor_columns = {
                row["name"]
                for row in connection.execute("PRAGMA table_info(professors)")
            }
            if "openai_api_key_encrypted" not in professor_columns:
                connection.execute(
                    "ALTER TABLE professors ADD COLUMN openai_api_key_encrypted TEXT"
                )
            if "openai_api_key_last_four" not in professor_columns:
                connection.execute(
                    "ALTER TABLE professors ADD COLUMN openai_api_key_last_four TEXT"
                )
            course_columns = {
                row["name"]
                for row in connection.execute("PRAGMA table_info(courses)")
            }
            if "provider" not in course_columns:
                connection.execute(
                    "ALTER TABLE courses ADD COLUMN provider TEXT NOT NULL DEFAULT 'anthropic'"
                )
            if "model" not in course_columns:
                connection.execute(
                    "ALTER TABLE courses ADD COLUMN model TEXT NOT NULL DEFAULT 'claude-sonnet-4-6'"
                )
            if "model_updated_at" not in course_columns:
                connection.execute(
                    "ALTER TABLE courses ADD COLUMN model_updated_at TEXT"
                )
            if "project_builder_enabled" not in course_columns:
                connection.execute(
                    """
                    ALTER TABLE courses
                    ADD COLUMN project_builder_enabled INTEGER NOT NULL DEFAULT 1
                    """
                )
            if "research_innovation_enabled" not in course_columns:
                connection.execute(
                    """
                    ALTER TABLE courses
                    ADD COLUMN research_innovation_enabled INTEGER NOT NULL DEFAULT 0
                    """
                )
            usage_columns = {
                row["name"]
                for row in connection.execute("PRAGMA table_info(usage_events)")
            }
            if "event_type" not in usage_columns:
                connection.execute(
                    """
                    ALTER TABLE usage_events
                    ADD COLUMN event_type TEXT NOT NULL DEFAULT 'student_question'
                    """
                )
            if "provider" not in usage_columns:
                connection.execute(
                    """
                    ALTER TABLE usage_events
                    ADD COLUMN provider TEXT NOT NULL DEFAULT 'anthropic'
                    """
                )
            document_columns = {
                row["name"]
                for row in connection.execute("PRAGMA table_info(documents)")
            }
            if "extracted_chars" not in document_columns:
                connection.execute(
                    """
                    ALTER TABLE documents
                    ADD COLUMN extracted_chars INTEGER NOT NULL DEFAULT 0
                    """
                )
        try:
            self.db_path.chmod(0o600)
        except OSError:
            pass

    def create_session(
        self,
        role: str,
        professor_id: Optional[str] = None,
        lifetime_hours: int = 12,
    ) -> str:
        if role not in {"admin", "professor"}:
            raise PilotValidationError("Invalid session role.")
        if role == "professor" and not self.get_professor(professor_id or ""):
            raise PilotValidationError("Professor account not found.")
        if role == "admin":
            professor_id = None
        lifetime_hours = min(max(int(lifetime_hours), 1), 168)
        now = datetime.now(timezone.utc)
        token = secrets.token_urlsafe(32)
        with self._connect() as connection:
            connection.execute(
                "DELETE FROM sessions WHERE expires_at <= ?", (now.isoformat(),)
            )
            connection.execute(
                """
                INSERT INTO sessions
                    (token_hash, role, professor_id, expires_at, created_at)
                VALUES (?, ?, ?, ?, ?)
                """,
                (
                    hash_token(token),
                    role,
                    professor_id,
                    (now + timedelta(hours=lifetime_hours)).isoformat(),
                    now.isoformat(),
                ),
            )
        return token

    def session_details(self, token: str) -> Optional[Dict]:
        if not token:
            return None
        now = utc_now()
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT token_hash, role, professor_id, expires_at, created_at
                FROM sessions WHERE token_hash = ?
                """,
                (hash_token(token),),
            ).fetchone()
            if row and row["expires_at"] <= now:
                connection.execute(
                    "DELETE FROM sessions WHERE token_hash = ?",
                    (hash_token(token),),
                )
                return None
        return self._row_dict(row)

    def delete_session(self, token: str) -> None:
        if not token:
            return
        with self._connect() as connection:
            connection.execute(
                "DELETE FROM sessions WHERE token_hash = ?", (hash_token(token),)
            )

    @staticmethod
    def _row_dict(row: Optional[sqlite3.Row]) -> Optional[Dict]:
        return dict(row) if row is not None else None

    def professor_count(self) -> int:
        with self._connect() as connection:
            row = connection.execute(
                "SELECT COUNT(*) AS count FROM professors WHERE is_active = 1"
            ).fetchone()
        return int(row["count"])

    def _reserved_pilot_slots(self, connection: sqlite3.Connection) -> int:
        now = utc_now()
        active = connection.execute(
            "SELECT COUNT(*) AS count FROM professors WHERE is_active = 1"
        ).fetchone()["count"]
        pending = connection.execute(
            """
            SELECT COUNT(*) AS count
            FROM invitations
            WHERE used_at IS NULL AND expires_at > ?
            """,
            (now,),
        ).fetchone()["count"]
        return int(active) + int(pending)

    def create_invitation(
        self, email: str, name: str, expires_hours: int = 72
    ) -> Tuple[Dict, str]:
        normalized_email = normalize_email(email)
        clean_name = _clean_text(name, "Name", 120)
        expires_hours = min(max(int(expires_hours), 1), 168)
        token = secrets.token_urlsafe(32)
        now = datetime.now(timezone.utc)
        invitation = {
            "id": str(secrets.token_hex(16)),
            "email": normalized_email,
            "name": clean_name,
            "token_hash": hash_token(token),
            "expires_at": (now + timedelta(hours=expires_hours)).isoformat(),
            "used_at": None,
            "created_at": now.isoformat(),
        }

        with self._connect() as connection:
            connection.execute("BEGIN IMMEDIATE")
            existing = connection.execute(
                "SELECT id FROM professors WHERE email = ?", (normalized_email,)
            ).fetchone()
            if existing:
                raise PilotValidationError("That professor already has an account.")
            # Renewing an invitation for the same person should not consume a
            # second pilot slot.
            connection.execute(
                """
                UPDATE invitations
                SET expires_at = ?
                WHERE email = ? AND used_at IS NULL
                """,
                (now.isoformat(), normalized_email),
            )
            if self._reserved_pilot_slots(connection) >= self.max_professors:
                raise PilotValidationError(
                    f"The {self.max_professors}-professor pilot cohort is full."
                )
            connection.execute(
                """
                INSERT INTO invitations
                    (id, email, name, token_hash, expires_at, used_at, created_at)
                VALUES
                    (:id, :email, :name, :token_hash, :expires_at, :used_at, :created_at)
                """,
                invitation,
            )
        public_invitation = {k: invitation[k] for k in (
            "id", "email", "name", "expires_at", "created_at"
        )}
        return public_invitation, token

    def invitation_details(self, token: str) -> Optional[Dict]:
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT id, email, name, expires_at, used_at, created_at
                FROM invitations WHERE token_hash = ?
                """,
                (hash_token(token),),
            ).fetchone()
        invitation = self._row_dict(row)
        if not invitation or invitation["used_at"]:
            return None
        if invitation["expires_at"] <= utc_now():
            return None
        return invitation

    def accept_invitation(self, token: str, password: str, name: str = "") -> Dict:
        password_hash = hash_password(password)
        now = utc_now()
        with self._connect() as connection:
            connection.execute("BEGIN IMMEDIATE")
            invitation = connection.execute(
                "SELECT * FROM invitations WHERE token_hash = ?",
                (hash_token(token),),
            ).fetchone()
            if not invitation or invitation["used_at"]:
                raise PilotValidationError("This invitation is invalid or has already been used.")
            if invitation["expires_at"] <= now:
                raise PilotValidationError("This invitation has expired.")
            active_count = connection.execute(
                "SELECT COUNT(*) AS count FROM professors WHERE is_active = 1"
            ).fetchone()["count"]
            if int(active_count) >= self.max_professors:
                raise PilotValidationError("The pilot cohort is full.")

            professor_id = "p_" + secrets.token_urlsafe(16)
            professor_name = _clean_text(
                name or invitation["name"], "Name", 120
            )
            try:
                connection.execute(
                    """
                    INSERT INTO professors
                        (id, email, name, password_hash, is_active, created_at)
                    VALUES (?, ?, ?, ?, 1, ?)
                    """,
                    (
                        professor_id,
                        invitation["email"],
                        professor_name,
                        password_hash,
                        now,
                    ),
                )
            except sqlite3.IntegrityError as exc:
                raise PilotValidationError("That professor already has an account.") from exc
            connection.execute(
                "UPDATE invitations SET used_at = ? WHERE id = ?",
                (now, invitation["id"]),
            )
        return self.get_professor(professor_id)

    def authenticate_professor(self, email: str, password: str) -> Optional[Dict]:
        try:
            normalized_email = normalize_email(email)
        except PilotValidationError:
            return None
        with self._connect() as connection:
            row = connection.execute(
                "SELECT * FROM professors WHERE email = ? AND is_active = 1",
                (normalized_email,),
            ).fetchone()
        stored_hash = row["password_hash"] if row else _DUMMY_PASSWORD_HASH
        if not verify_password(password, stored_hash) or not row:
            return None
        return self._public_professor(dict(row))

    @staticmethod
    def _public_professor(professor: Dict) -> Dict:
        anthropic_key = bool(professor.get("api_key_encrypted"))
        openai_key = bool(professor.get("openai_api_key_encrypted"))
        return {
            "id": professor["id"],
            "email": professor["email"],
            "name": professor["name"],
            # The two fields below remain as Anthropic aliases for older clients.
            "has_api_key": anthropic_key,
            "api_key_last_four": professor.get("api_key_last_four") or "",
            "api_keys": {
                "anthropic": {
                    "has_key": anthropic_key,
                    "last_four": professor.get("api_key_last_four") or "",
                },
                "openai": {
                    "has_key": openai_key,
                    "last_four": professor.get("openai_api_key_last_four") or "",
                },
            },
            "is_active": bool(professor.get("is_active", 1)),
            "created_at": professor["created_at"],
        }

    def get_professor(self, professor_id: str) -> Optional[Dict]:
        with self._connect() as connection:
            row = connection.execute(
                "SELECT * FROM professors WHERE id = ? AND is_active = 1",
                (professor_id,),
            ).fetchone()
        return self._public_professor(dict(row)) if row else None

    @staticmethod
    def _api_key_columns(provider: str) -> Tuple[str, str]:
        try:
            normalized = normalize_provider(provider)
        except ValueError as exc:
            raise PilotValidationError(str(exc)) from exc
        if normalized == "anthropic":
            return "api_key_encrypted", "api_key_last_four"
        return "openai_api_key_encrypted", "openai_api_key_last_four"

    @staticmethod
    def _validate_api_key(provider: str, api_key: str) -> str:
        key = (api_key or "").strip()
        if provider == "anthropic":
            valid = key.startswith("sk-ant-") and len(key) >= 30
            message = "Enter a valid Anthropic API key."
        else:
            valid = key.startswith("sk-") and len(key) >= 20
            message = "Enter a valid OpenAI API key."
        if not valid:
            raise PilotValidationError(message)
        return key

    def set_professor_api_key(
        self,
        professor_id: str,
        api_key: str,
        provider: str = "anthropic",
    ) -> Dict:
        try:
            normalized = normalize_provider(provider)
        except ValueError as exc:
            raise PilotValidationError(str(exc)) from exc
        encrypted_column, last_four_column = self._api_key_columns(normalized)
        key = self._validate_api_key(normalized, api_key)
        encrypted = self.fernet.encrypt(key.encode("utf-8")).decode("ascii")
        with self._connect() as connection:
            result = connection.execute(
                f"""
                UPDATE professors SET {encrypted_column} = ?, {last_four_column} = ?
                WHERE id = ? AND is_active = 1
                """,  # Column names come only from _api_key_columns.
                (encrypted, key[-4:], professor_id),
            )
            if result.rowcount != 1:
                raise PilotValidationError("Professor account not found.")
        return self.get_professor(professor_id)

    def delete_professor_api_key(
        self, professor_id: str, provider: str = "anthropic"
    ) -> Dict:
        encrypted_column, last_four_column = self._api_key_columns(provider)
        with self._connect() as connection:
            connection.execute(
                f"""
                UPDATE professors SET {encrypted_column} = NULL, {last_four_column} = NULL
                WHERE id = ?
                """,  # Column names come only from _api_key_columns.
                (professor_id,),
            )
        return self.get_professor(professor_id)

    def decrypted_api_key(
        self, professor_id: str, provider: str = "anthropic"
    ) -> Optional[str]:
        encrypted_column, _ = self._api_key_columns(provider)
        with self._connect() as connection:
            row = connection.execute(
                f"""
                SELECT {encrypted_column} AS encrypted_key FROM professors
                WHERE id = ? AND is_active = 1
                """,  # Column name comes only from _api_key_columns.
                (professor_id,),
            ).fetchone()
        if not row or not row["encrypted_key"]:
            return None
        try:
            return self.fernet.decrypt(
                row["encrypted_key"].encode("ascii")
            ).decode("utf-8")
        except InvalidToken as exc:
            raise PilotConfigurationError(
                "The stored API key cannot be decrypted with ATLAS_ENCRYPTION_KEY."
            ) from exc

    def create_course(
        self,
        owner_id: str,
        name: str,
        code: str,
        term: str,
        section: str = "",
        campus: str = "Arlington",
        monthly_question_limit: int = DEFAULT_MONTHLY_QUESTION_LIMIT,
        provider: str = "anthropic",
        model: str = DEFAULT_MODEL_BY_PROVIDER["anthropic"],
    ) -> Dict:
        if not self.get_professor(owner_id):
            raise PilotValidationError("Professor account not found.")
        with self._connect() as connection:
            course_count = connection.execute(
                "SELECT COUNT(*) AS count FROM courses WHERE owner_id = ?",
                (owner_id,),
            ).fetchone()["count"]
        if int(course_count) >= MAX_COURSES_PER_PROFESSOR:
            raise PilotValidationError(
                f"Each professor can create at most {MAX_COURSES_PER_PROFESSOR} courses."
            )
        clean_name = _clean_text(name, "Course name", 160)
        clean_code = _clean_text(code, "Course code", 40).upper()
        clean_term = _clean_text(term, "Term", 80)
        clean_section = _clean_text(section, "Section", 80, required=False)
        clean_campus = _clean_text(campus, "Campus", 80)
        try:
            clean_provider, clean_model = validate_provider_model(provider, model)
        except ValueError as exc:
            raise PilotValidationError(str(exc)) from exc
        limit = min(max(int(monthly_question_limit), 1), 5000)
        course_id = "c_" + secrets.token_urlsafe(20)
        now = utc_now()
        with self._connect() as connection:
            connection.execute(
                """
                INSERT INTO courses
                    (id, owner_id, name, code, section, term, campus, status,
                     monthly_question_limit, concept_map_json, provider, model,
                     model_updated_at, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, 'draft', ?, '{}', ?, ?, ?, ?, ?)
                """,
                (
                    course_id,
                    owner_id,
                    clean_name,
                    clean_code,
                    clean_section,
                    clean_term,
                    clean_campus,
                    limit,
                    clean_provider,
                    clean_model,
                    now,
                    now,
                    now,
                ),
            )
        (self.courses_dir / course_id / "originals").mkdir(parents=True, exist_ok=True)
        (self.courses_dir / course_id / "extracted").mkdir(parents=True, exist_ok=True)
        return self.get_course(course_id)

    def get_course(self, course_id: str) -> Optional[Dict]:
        with self._connect() as connection:
            row = connection.execute(
                "SELECT * FROM courses WHERE id = ?", (course_id,)
            ).fetchone()
        return self._row_dict(row)

    def list_courses_for_professor(self, professor_id: str) -> List[Dict]:
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT * FROM courses
                WHERE owner_id = ?
                ORDER BY created_at DESC
                """,
                (professor_id,),
            ).fetchall()
        return [dict(row) for row in rows]

    def list_all_courses(self) -> List[Dict]:
        with self._connect() as connection:
            rows = connection.execute(
                "SELECT * FROM courses ORDER BY created_at DESC"
            ).fetchall()
        return [dict(row) for row in rows]

    def set_course_status(self, course_id: str, owner_id: str, status: str) -> Dict:
        if status not in {"draft", "published", "archived"}:
            raise PilotValidationError("Invalid course status.")
        with self._connect() as connection:
            result = connection.execute(
                """
                UPDATE courses SET status = ?, updated_at = ?
                WHERE id = ? AND owner_id = ?
                """,
                (status, utc_now(), course_id, owner_id),
            )
            if result.rowcount != 1:
                raise PilotValidationError("Course not found.")
        return self.get_course(course_id)

    def set_course_features(
        self,
        course_id: str,
        owner_id: str,
        project_builder_enabled: bool,
        research_innovation_enabled: bool,
    ) -> Dict:
        """Update student guide visibility without changing any course content."""
        with self._connect() as connection:
            result = connection.execute(
                """
                UPDATE courses
                SET project_builder_enabled = ?,
                    research_innovation_enabled = ?,
                    updated_at = ?
                WHERE id = ? AND owner_id = ?
                """,
                (
                    int(bool(project_builder_enabled)),
                    int(bool(research_innovation_enabled)),
                    utc_now(),
                    course_id,
                    owner_id,
                ),
            )
            if result.rowcount != 1:
                raise PilotValidationError("Course not found.")
        return self.get_course(course_id)

    def legacy_course_features(self, course_id: str) -> Dict:
        """Return persistent guide settings for a repository-backed course."""
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT project_builder_enabled, research_innovation_enabled,
                       updated_at
                FROM legacy_course_features
                WHERE course_id = ?
                """,
                (course_id,),
            ).fetchone()
        if not row:
            return {
                "project_builder_enabled": True,
                "research_innovation_enabled": False,
                "updated_at": None,
            }
        return {
            "project_builder_enabled": bool(row["project_builder_enabled"]),
            "research_innovation_enabled": bool(
                row["research_innovation_enabled"]
            ),
            "updated_at": row["updated_at"],
        }

    def set_legacy_course_features(
        self,
        course_id: str,
        project_builder_enabled: bool,
        research_innovation_enabled: bool,
    ) -> Dict:
        """Persist guide settings for a repository-backed legacy course."""
        now = utc_now()
        with self._connect() as connection:
            connection.execute(
                """
                INSERT INTO legacy_course_features
                    (course_id, project_builder_enabled,
                     research_innovation_enabled, updated_at)
                VALUES (?, ?, ?, ?)
                ON CONFLICT(course_id) DO UPDATE SET
                    project_builder_enabled = excluded.project_builder_enabled,
                    research_innovation_enabled = excluded.research_innovation_enabled,
                    updated_at = excluded.updated_at
                """,
                (
                    course_id,
                    int(bool(project_builder_enabled)),
                    int(bool(research_innovation_enabled)),
                    now,
                ),
            )
        return self.legacy_course_features(course_id)

    def set_course_model(
        self,
        course_id: str,
        owner_id: str,
        provider: str,
        model: str,
    ) -> Dict:
        try:
            clean_provider, clean_model = validate_provider_model(provider, model)
        except ValueError as exc:
            raise PilotValidationError(str(exc)) from exc
        if not self.decrypted_api_key(owner_id, clean_provider):
            provider_name = "Anthropic" if clean_provider == "anthropic" else "OpenAI"
            raise PilotValidationError(
                f"Add your {provider_name} API key before selecting that provider."
            )

        now = utc_now()
        with self._connect() as connection:
            connection.execute("BEGIN IMMEDIATE")
            course = connection.execute(
                "SELECT * FROM courses WHERE id = ? AND owner_id = ?",
                (course_id, owner_id),
            ).fetchone()
            if not course:
                raise PilotValidationError("Course not found.")
            if course["provider"] == clean_provider and course["model"] == clean_model:
                return dict(course)
            connection.execute(
                """
                UPDATE courses
                SET provider = ?, model = ?, model_updated_at = ?, updated_at = ?
                WHERE id = ? AND owner_id = ?
                """,
                (clean_provider, clean_model, now, now, course_id, owner_id),
            )
            connection.execute(
                """
                INSERT INTO course_model_events
                    (id, course_id, professor_id, old_provider, old_model,
                     new_provider, new_model, created_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    "m_" + secrets.token_urlsafe(16),
                    course_id,
                    owner_id,
                    course["provider"],
                    course["model"],
                    clean_provider,
                    clean_model,
                    now,
                ),
            )
        return self.get_course(course_id)

    def course_model_history(
        self, course_id: str, owner_id: str, limit: int = 20
    ) -> List[Dict]:
        course = self.get_course(course_id)
        if not course or course["owner_id"] != owner_id:
            raise PilotValidationError("Course not found.")
        safe_limit = min(max(int(limit), 1), 100)
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT old_provider, old_model, new_provider, new_model, created_at
                FROM course_model_events
                WHERE course_id = ? AND professor_id = ?
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (course_id, owner_id, safe_limit),
            ).fetchall()
        return [dict(row) for row in rows]

    def set_concept_map(self, course_id: str, owner_id: str, concept_map: Dict) -> Dict:
        serialized = json.dumps(concept_map, ensure_ascii=False)
        with self._connect() as connection:
            result = connection.execute(
                """
                UPDATE courses SET concept_map_json = ?, updated_at = ?
                WHERE id = ? AND owner_id = ?
                """,
                (serialized, utc_now(), course_id, owner_id),
            )
            if result.rowcount != 1:
                raise PilotValidationError("Course not found.")
        return self.get_course(course_id)

    def save_document(
        self,
        course_id: str,
        owner_id: str,
        filename: str,
        document_type: str,
        content: bytes,
        extracted_text: str,
    ) -> Dict:
        if document_type not in {"syllabus", "transcript"}:
            raise PilotValidationError("Invalid document type.")
        course = self.get_course(course_id)
        if not course or course["owner_id"] != owner_id:
            raise PilotValidationError("Course not found.")
        if not extracted_text.strip():
            raise PilotValidationError("No readable text was found in that document.")

        original_name = safe_filename(filename)
        extension = Path(original_name).suffix.lower()
        with self._connect() as connection:
            totals = connection.execute(
                """
                SELECT COUNT(*) AS document_count,
                       COALESCE(SUM(byte_size), 0) AS total_bytes,
                       COALESCE(SUM(extracted_chars), 0) AS extracted_chars
                FROM documents WHERE course_id = ?
                """,
                (course_id,),
            ).fetchone()
            replaced = connection.execute(
                """
                SELECT COUNT(*) AS document_count,
                       COALESCE(SUM(byte_size), 0) AS total_bytes,
                       COALESCE(SUM(extracted_chars), 0) AS extracted_chars
                FROM documents
                WHERE course_id = ? AND document_type = 'syllabus'
                """,
                (course_id,),
            ).fetchone() if document_type == "syllabus" else None

        current_count = int(totals["document_count"])
        current_bytes = int(totals["total_bytes"])
        replaced_count = int(replaced["document_count"]) if replaced else 0
        replaced_bytes = int(replaced["total_bytes"]) if replaced else 0
        current_extracted_chars = int(totals["extracted_chars"])
        replaced_extracted_chars = int(replaced["extracted_chars"]) if replaced else 0
        projected_count = current_count - replaced_count + 1
        projected_bytes = current_bytes - replaced_bytes + len(content)
        projected_extracted_chars = (
            current_extracted_chars - replaced_extracted_chars + len(extracted_text)
        )
        if projected_count > MAX_COURSE_DOCUMENTS:
            raise PilotValidationError(
                f"A course can contain at most {MAX_COURSE_DOCUMENTS} documents."
            )
        if projected_bytes > MAX_COURSE_DOCUMENT_BYTES:
            raise PilotValidationError(
                "Course documents can use at most 250 MB of storage."
            )
        if projected_extracted_chars > MAX_COURSE_EXTRACTED_CHARS:
            raise PilotValidationError(
                "The course contains too much extracted text. Remove or shorten materials."
            )

        document_id = "d_" + secrets.token_urlsafe(16)
        course_dir = (self.courses_dir / course_id).resolve()
        originals_dir = course_dir / "originals"
        extracted_dir = course_dir / "extracted"
        originals_dir.mkdir(parents=True, exist_ok=True)
        extracted_dir.mkdir(parents=True, exist_ok=True)
        try:
            course_dir.chmod(0o700)
            originals_dir.chmod(0o700)
            extracted_dir.chmod(0o700)
        except OSError:
            pass
        stored_path = originals_dir / f"{document_id}{extension}"
        extracted_path = extracted_dir / f"{document_id}.txt"
        stored_path.write_bytes(content)
        extracted_path.write_text(extracted_text, encoding="utf-8")
        try:
            stored_path.chmod(0o600)
            extracted_path.chmod(0o600)
        except OSError:
            pass

        now = utc_now()
        record = {
            "id": document_id,
            "course_id": course_id,
            "filename": original_name,
            "document_type": document_type,
            "stored_path": str(stored_path.relative_to(self.data_dir)),
            "extracted_path": str(extracted_path.relative_to(self.data_dir)),
            "byte_size": len(content),
            "extracted_chars": len(extracted_text),
            "sha256": hashlib.sha256(content).hexdigest(),
            "uploaded_at": now,
        }
        with self._connect() as connection:
            connection.execute("BEGIN IMMEDIATE")
            if document_type == "syllabus":
                old_rows = connection.execute(
                    """
                    SELECT id, stored_path, extracted_path FROM documents
                    WHERE course_id = ? AND document_type = 'syllabus'
                    """,
                    (course_id,),
                ).fetchall()
                connection.execute(
                    """
                    DELETE FROM documents
                    WHERE course_id = ? AND document_type = 'syllabus'
                    """,
                    (course_id,),
                )
                for old in old_rows:
                    for relative_path in (old["stored_path"], old["extracted_path"]):
                        old_path = self.data_dir / relative_path
                        if old_path.exists():
                            old_path.unlink()
            connection.execute(
                """
                INSERT INTO documents
                    (id, course_id, filename, document_type, stored_path,
                     extracted_path, byte_size, extracted_chars, sha256, uploaded_at)
                VALUES
                    (:id, :course_id, :filename, :document_type, :stored_path,
                     :extracted_path, :byte_size, :extracted_chars, :sha256, :uploaded_at)
                """,
                record,
            )
            connection.execute(
                "UPDATE courses SET updated_at = ? WHERE id = ?",
                (now, course_id),
            )
        return record

    def list_documents(self, course_id: str, owner_id: Optional[str] = None) -> List[Dict]:
        course = self.get_course(course_id)
        if not course or (owner_id and course["owner_id"] != owner_id):
            raise PilotValidationError("Course not found.")
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT id, course_id, filename, document_type, byte_size,
                       extracted_chars, sha256, uploaded_at
                FROM documents WHERE course_id = ?
                ORDER BY uploaded_at DESC
                """,
                (course_id,),
            ).fetchall()
        return [dict(row) for row in rows]

    def load_course_materials(self, course_id: str) -> Tuple[str, Dict[str, str], Dict]:
        course = self.get_course(course_id)
        if not course:
            raise PilotValidationError("Course not found.")
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT filename, document_type, extracted_path
                FROM documents WHERE course_id = ?
                ORDER BY uploaded_at ASC
                """,
                (course_id,),
            ).fetchall()

        syllabus = ""
        sources: Dict[str, str] = {}
        for row in rows:
            path = (self.data_dir / row["extracted_path"]).resolve()
            try:
                path.relative_to(self.data_dir)
            except ValueError as exc:
                raise PilotConfigurationError("Unsafe stored document path.") from exc
            if not path.exists():
                continue
            text = path.read_text(encoding="utf-8")
            if row["document_type"] == "syllabus":
                syllabus = text
            else:
                display_name = row["filename"]
                if display_name in sources:
                    display_name = f"{Path(display_name).stem}-{len(sources) + 1}{Path(display_name).suffix}"
                sources[display_name] = text
        try:
            concept_map = json.loads(course["concept_map_json"] or "{}")
        except json.JSONDecodeError:
            concept_map = {}
        return syllabus, sources, concept_map

    def delete_document(self, document_id: str, course_id: str, owner_id: str) -> None:
        course = self.get_course(course_id)
        if not course or course["owner_id"] != owner_id:
            raise PilotValidationError("Course not found.")
        with self._connect() as connection:
            connection.execute("BEGIN IMMEDIATE")
            row = connection.execute(
                """
                SELECT stored_path, extracted_path FROM documents
                WHERE id = ? AND course_id = ?
                """,
                (document_id, course_id),
            ).fetchone()
            if not row:
                raise PilotValidationError("Document not found.")
            connection.execute(
                "DELETE FROM documents WHERE id = ? AND course_id = ?",
                (document_id, course_id),
            )
        for relative_path in (row["stored_path"], row["extracted_path"]):
            path = self.data_dir / relative_path
            if path.exists():
                path.unlink()

    def monthly_usage(self, course_id: str) -> Dict:
        month_start = datetime.now(timezone.utc).replace(
            day=1, hour=0, minute=0, second=0, microsecond=0
        ).isoformat()
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT COALESCE(
                           SUM(CASE WHEN event_type = 'student_question' THEN 1 ELSE 0 END),
                           0
                       ) AS questions,
                       COALESCE(SUM(CASE WHEN event_type = 'concept_map' THEN 1 ELSE 0 END), 0)
                           AS concept_map_generations,
                       COALESCE(SUM(input_tokens), 0) AS input_tokens,
                       COALESCE(SUM(output_tokens), 0) AS output_tokens
                FROM usage_events
                WHERE course_id = ? AND created_at >= ?
                """,
                (course_id, month_start),
            ).fetchone()
            breakdown = connection.execute(
                """
                SELECT provider, model,
                       COALESCE(
                           SUM(CASE WHEN event_type = 'student_question' THEN 1 ELSE 0 END),
                           0
                       ) AS questions,
                       COALESCE(SUM(CASE WHEN event_type = 'concept_map' THEN 1 ELSE 0 END), 0)
                           AS concept_map_generations,
                       COALESCE(SUM(input_tokens), 0) AS input_tokens,
                       COALESCE(SUM(output_tokens), 0) AS output_tokens
                FROM usage_events
                WHERE course_id = ? AND created_at >= ?
                GROUP BY provider, model
                ORDER BY provider, model
                """,
                (course_id, month_start),
            ).fetchall()
        return {
            "questions": int(row["questions"]),
            "concept_map_generations": int(row["concept_map_generations"]),
            "input_tokens": int(row["input_tokens"]),
            "output_tokens": int(row["output_tokens"]),
            "by_model": [
                {
                    "provider": item["provider"],
                    "model": item["model"],
                    "questions": int(item["questions"]),
                    "concept_map_generations": int(item["concept_map_generations"]),
                    "input_tokens": int(item["input_tokens"]),
                    "output_tokens": int(item["output_tokens"]),
                }
                for item in breakdown
            ],
        }

    def remaining_questions(self, course_id: str) -> int:
        course = self.get_course(course_id)
        if not course:
            return 0
        used = self.monthly_usage(course_id)["questions"]
        return max(0, int(course["monthly_question_limit"]) - used)

    def record_usage(
        self,
        course_id: str,
        professor_id: str,
        session_id: str,
        model: str,
        input_tokens: int,
        output_tokens: int,
        event_type: str = "student_question",
        provider: str = "anthropic",
    ) -> None:
        if event_type not in {"student_question", "concept_map"}:
            raise PilotValidationError("Invalid usage event type.")
        try:
            clean_provider = normalize_provider(provider)
        except ValueError as exc:
            raise PilotValidationError(str(exc)) from exc
        with self._connect() as connection:
            connection.execute(
                """
                INSERT INTO usage_events
                    (id, course_id, professor_id, session_id, provider, model, event_type,
                     input_tokens, output_tokens, created_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    "u_" + secrets.token_urlsafe(16),
                    course_id,
                    professor_id,
                    session_id,
                    clean_provider,
                    model,
                    event_type,
                    max(0, int(input_tokens)),
                    max(0, int(output_tokens)),
                    utc_now(),
                ),
            )

    def record_feedback(
        self,
        course_id: str,
        session_id: str,
        rating: str,
        comment: str = "",
    ) -> None:
        if rating not in {"up", "down"}:
            raise PilotValidationError("Invalid feedback rating.")
        clean_comment = (comment or "").strip()[:1000]
        with self._connect() as connection:
            connection.execute(
                """
                INSERT INTO pilot_feedback
                    (id, course_id, session_id, rating, comment, created_at)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    "f_" + secrets.token_urlsafe(16),
                    course_id,
                    session_id,
                    rating,
                    clean_comment,
                    utc_now(),
                ),
            )

    def admin_summary(self) -> Dict:
        with self._connect() as connection:
            professors = connection.execute(
                """
                SELECT p.id, p.email, p.name, p.api_key_last_four,
                       p.openai_api_key_last_four, p.is_active,
                       p.created_at, COUNT(DISTINCT c.id) AS course_count
                FROM professors p
                LEFT JOIN courses c ON c.owner_id = p.id
                GROUP BY p.id
                ORDER BY p.created_at ASC
                """
            ).fetchall()
            pending = connection.execute(
                """
                SELECT id, email, name, expires_at, created_at
                FROM invitations
                WHERE used_at IS NULL AND expires_at > ?
                ORDER BY created_at DESC
                """,
                (utc_now(),),
            ).fetchall()
        return {
            "max_professors": self.max_professors,
            "active_professors": self.professor_count(),
            "professors": [dict(row) for row in professors],
            "pending_invitations": [dict(row) for row in pending],
        }


def build_store_from_environment() -> PilotStore:
    data_dir = os.getenv("ATLAS_DATA_DIR", "").strip()
    encryption_key = os.getenv("ATLAS_ENCRYPTION_KEY", "").strip()
    if not data_dir:
        raise PilotConfigurationError("ATLAS_DATA_DIR is required when the pilot is enabled.")
    if not encryption_key:
        raise PilotConfigurationError(
            "ATLAS_ENCRYPTION_KEY is required when the pilot is enabled."
        )
    try:
        max_professors = int(
            os.getenv("ATLAS_MAX_PILOT_PROFESSORS", str(DEFAULT_MAX_PROFESSORS))
        )
    except ValueError as exc:
        raise PilotConfigurationError(
            "ATLAS_MAX_PILOT_PROFESSORS must be an integer."
        ) from exc
    store = PilotStore(Path(data_dir), encryption_key, max_professors)
    store.initialize()
    return store
